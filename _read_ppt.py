import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'])
    from pptx import Presentation

pptx_path = r'C:\Users\sudpy\Downloads\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx'
prs = Presentation(pptx_path)
print(f'Total slides: {len(prs.slides)}')
print()
for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        # Also note image placeholders
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            texts.append(f'[IMAGE: {shape.name}]')
        if 'placeholder' in shape.name.lower() and not shape.has_text_frame:
            texts.append(f'[PLACEHOLDER: {shape.name}]')
    print(f'--- SLIDE {i+1} [{layout}] ---')
    for t in texts[:30]:
        print(f'  {repr(t)}')
    print()
